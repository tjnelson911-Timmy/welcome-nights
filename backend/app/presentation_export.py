"""
Welcome Nights Presentation Export Service

This module provides export functionality for presentations to PPTX and PDF formats.
Uses the AV3 Culture Night template as a base for generating presentations.
"""

import io
import os
import tempfile
import urllib.request
from typing import List, Dict, Any, Optional
from datetime import datetime
from pathlib import Path
from sqlalchemy.orm import Session
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from . import models, crud

# Mapbox token for static map images
MAPBOX_TOKEN = "pk.eyJ1IjoidGpuZWxzb245MTEiLCJhIjoiY21rZ2hvd2h6MDc1bDNkb256c2ZpZzJ5ZSJ9.GjbF9IEBFXgJl-unUW4hoQ"


# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Template path
TEMPLATE_DIR = Path(__file__).parent.parent / "templates"
AV3_TEMPLATE = TEMPLATE_DIR / "AV3 Culture Night Presentation - Salem (1)_20260128_020248.pptx"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def export_to_pptx(
    db: Session,
    presentation_id: int,
    output_path: Optional[str] = None
) -> bytes:
    """
    Export a presentation to PPTX format using the AV3 template.

    Args:
        db: Database session
        presentation_id: ID of the presentation to export
        output_path: Optional path to save the file (if None, returns bytes)

    Returns:
        Bytes of the PPTX file
    """
    # Get presentation data
    presentation = crud.get_presentation(db, presentation_id)
    if not presentation:
        raise ValueError(f"Presentation {presentation_id} not found")

    brand = crud.get_brand(db, presentation.brand_id)
    facility = crud.get_facility(db, presentation.facility_id)
    slides = crud.get_presentation_slides(db, presentation_id)

    # Get facility logo from assigned logo_asset_id
    facility_logo_path = None
    if facility.logo_asset_id:
        asset = crud.get_asset(db, facility.logo_asset_id)
        if asset and asset.url:
            if asset.url.startswith('/uploads/'):
                facility_logo_path = Path(__file__).parent.parent / asset.url.lstrip('/')
            elif asset.url.startswith('uploads/'):
                facility_logo_path = Path(__file__).parent.parent / asset.url

    # Check if template exists
    if AV3_TEMPLATE.exists():
        return export_using_template(
            db, presentation, brand, facility, slides, facility_logo_path, output_path
        )
    else:
        # Fallback to generated slides
        return export_generated_slides(
            brand, facility, slides, output_path
        )


def export_using_template(
    db: Session,
    presentation: models.Presentation,
    brand: models.Brand,
    facility: models.Facility,
    slides: List[models.PresentationSlide],
    facility_logo_path: Optional[Path] = None,
    output_path: Optional[str] = None
) -> bytes:
    """
    Export presentation using the AV3 template as a base.
    Replaces facility logo on title and end slides.
    Uses reusable content from database for history, footprint, regions, culture.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Load the template
    pptx = PptxPresentation(str(AV3_TEMPLATE))

    # Load reusable content for the brand
    content_items = crud.get_content(db, brand.id)
    content_map = {item.content_key: item for item in content_items}

    # Build a map of slide types to their indices in our slides data
    slide_data_map = {s.slide_type: s for s in slides}

    # Process each template slide
    for slide_idx, slide in enumerate(pptx.slides):
        # Replace logo on slide 1 (title) and slide 25 (end)
        if facility_logo_path and facility_logo_path.exists():
            if slide_idx == 0 or slide_idx == len(list(pptx.slides)) - 1:
                replace_slide_logo(slide, facility_logo_path)

        # Update text content with brand name and reusable content
        update_template_slide(db, slide, slide_idx, brand, facility, slide_data_map, content_map)

    # Save to bytes
    buffer = io.BytesIO()
    pptx.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.read()

    if output_path:
        with open(output_path, 'wb') as f:
            f.write(pptx_bytes)

    return pptx_bytes


def replace_slide_logo(slide, logo_path: Path):
    """
    Replace the main logo/picture on a slide with the facility logo.
    Finds the largest picture on the slide and replaces it.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    # Find the largest picture shape (likely the logo)
    largest_picture = None
    largest_area = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = shape.width * shape.height
            if area > largest_area:
                largest_area = area
                largest_picture = shape

    if largest_picture:
        # Store position and size
        left = largest_picture.left
        top = largest_picture.top
        width = largest_picture.width
        height = largest_picture.height

        # Remove old picture
        sp = largest_picture._element
        sp.getparent().remove(sp)

        # Add new picture at same position
        try:
            slide.shapes.add_picture(
                str(logo_path),
                left, top, width, height
            )
        except Exception as e:
            print(f"Error adding logo: {e}")


def update_template_slide(
    db: Session,
    slide,
    slide_idx: int,
    brand: models.Brand,
    facility: models.Facility,
    slide_data_map: Dict[str, models.PresentationSlide],
    content_map: Dict[str, models.ReusableContent] = None
):
    """
    Update a template slide by replacing placeholder text with actual content.
    Preserves all formatting, images, and design elements.
    Uses reusable content from database for specific slides.

    AV3 Template Structure:
    - Slide 5 (idx 4): History Timeline - dates and "The Cascadia Family"
    - Slide 6 (idx 5): Footprint Stats - numbers and labels
    - Slide 7 (idx 6): Regions - region names
    - Slide 8 (idx 7): Culture Intro - "WE ARE NOT CORPORATE!"
    - Slide 9 (idx 8): Culture Comparison - "The Common Way..."
    """
    content_map = content_map or {}

    # Define text replacements for brand customization
    replacements = {
        "Cascadia": brand.name,
        "cascadia": brand.name.lower(),
    }

    # Slide 5 (idx 4): History Timeline
    if slide_idx == 4 and 'history' in content_map:
        history = content_map['history']
        if history.content and 'items' in history.content:
            items = history.content['items']
            # Template has: "March 2015", "October 2021", "December 2025"
            template_dates = ["March 2015", "October 2021", "December 2025"]
            for i, item in enumerate(items[:len(template_dates)]):
                if i < len(template_dates):
                    replacements[template_dates[i]] = item.get('year', template_dates[i])

    # Slide 6 (idx 5): Footprint Stats with Map
    if slide_idx == 5:
        # Add map image from facility coordinates
        add_footprint_map(db, slide, brand)

        if 'footprint' in content_map:
            footprint = content_map['footprint']
            if footprint.content and 'stats' in footprint.content:
                stats = footprint.content['stats']
                # Template values: "48", "9", "1", "3,500", "4,000"
                # Template labels: "Skilled Nursing", "Senior Living ALF/ILF", "Home Health and Hospice", "Dedicated Employees", "Residents"
                template_values = ["48", "9", "1", "3,500", "4,000"]
                template_labels = ["Skilled Nursing", "Senior Living ALF/ILF", "Home Health and Hospice", "Dedicated Employees", "Residents"]
                for i, stat in enumerate(stats[:len(template_values)]):
                    if i < len(template_values):
                        replacements[template_values[i]] = str(stat.get('value', template_values[i]))
                    if i < len(template_labels):
                        replacements[template_labels[i]] = stat.get('label', template_labels[i])

    # Slide 7 (idx 6): Regions
    if slide_idx == 6 and 'regions' in content_map:
        regions = content_map['regions']
        if regions.content and 'regions' in regions.content:
            region_list = regions.content['regions']
            # Template regions
            template_regions = [
                "Columbia (OR, Western WA)",
                "Northern (Eastern WA, Northern ID, MT)",
                "3 Rivers (Central - ID)",
                "Envision (Southern – ID)",
                "Vincero (AZ)"
            ]
            for i, region in enumerate(region_list[:len(template_regions)]):
                if i < len(template_regions):
                    region_name = region.get('name', '')
                    if region.get('facilities'):
                        region_name += f" ({region['facilities']} facilities)"
                    replacements[template_regions[i]] = region_name

    # Process all shapes in the slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            update_text_frame(shape.text_frame, replacements)

        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    update_text_frame(cell.text_frame, replacements)

        # Handle grouped shapes
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                if sub_shape.has_text_frame:
                    update_text_frame(sub_shape.text_frame, replacements)


def add_footprint_map(db: Session, slide, brand: models.Brand):
    """
    Add a map image showing all facility locations to the footprint slide.
    """
    # Get all facilities for this brand with coordinates
    facilities = crud.get_facilities(db, brand_id=brand.id)
    facilities_with_coords = [f for f in facilities if f.latitude and f.longitude]

    if not facilities_with_coords:
        return  # No facilities with coordinates, skip map

    # Build Mapbox static map URL with markers
    markers = []
    for f in facilities_with_coords[:50]:  # Limit to 50 to avoid URL length issues
        markers.append(f"pin-s+0b7280({f.longitude},{f.latitude})")

    markers_str = ",".join(markers)
    map_url = f"https://api.mapbox.com/styles/v1/mapbox/light-v11/static/{markers_str}/auto/800x500@2x?access_token={MAPBOX_TOKEN}"

    try:
        # Download the map image
        with urllib.request.urlopen(map_url, timeout=30) as response:
            image_data = response.read()

        # Save to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            tmp_file.write(image_data)
            tmp_path = tmp_file.name

        # Add image to slide - position it on the right side of the slide
        # Slide is 13.333" x 7.5" (16:9)
        # Place map on the right half, with some padding
        left = Inches(6.5)
        top = Inches(1.5)
        width = Inches(6.3)
        height = Inches(4)

        slide.shapes.add_picture(tmp_path, left, top, width, height)

        # Clean up temp file
        os.unlink(tmp_path)

    except Exception as e:
        print(f"Error adding map to footprint slide: {e}")
        # Continue without the map if there's an error


def update_text_frame(text_frame, replacements: Dict[str, str]):
    """
    Update text in a text frame while preserving formatting.
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            new_text = original_text

            for old_text, new_text_value in replacements.items():
                if old_text in new_text:
                    new_text = new_text.replace(old_text, new_text_value)

            if new_text != original_text:
                run.text = new_text


def export_generated_slides(
    brand: models.Brand,
    facility: models.Facility,
    slides: List[models.PresentationSlide],
    output_path: Optional[str] = None
) -> bytes:
    """
    Fallback: Generate slides from scratch when template is not available.
    """
    # Create PPTX
    pptx = PptxPresentation()
    pptx.slide_width = SLIDE_WIDTH
    pptx.slide_height = SLIDE_HEIGHT

    # Get colors
    primary_color = hex_to_rgb(brand.primary_color or "#0b7280")
    secondary_color = hex_to_rgb(brand.secondary_color or "#065a67")

    # Generate slides
    for slide_data in slides:
        add_slide_to_pptx(
            pptx,
            slide_data,
            brand,
            facility,
            primary_color,
            secondary_color
        )

    # Save to bytes
    buffer = io.BytesIO()
    pptx.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.read()

    if output_path:
        with open(output_path, 'wb') as f:
            f.write(pptx_bytes)

    return pptx_bytes


def add_slide_to_pptx(
    pptx: PptxPresentation,
    slide_data: models.PresentationSlide,
    brand: models.Brand,
    facility: models.Facility,
    primary_color: RGBColor,
    secondary_color: RGBColor
):
    """Add a single slide to the PPTX based on slide type"""
    slide_type = slide_data.slide_type
    payload = slide_data.payload or {}

    # Use blank layout
    blank_layout = pptx.slide_layouts[6]  # Blank
    slide = pptx.slides.add_slide(blank_layout)

    if slide_type == "WelcomeIntro":
        create_welcome_pptx_slide(slide, payload, brand, facility, primary_color)
    elif slide_type == "RaffleBumper":
        create_raffle_pptx_slide(slide, payload, primary_color)
    elif slide_type == "HistoryBlock":
        create_content_pptx_slide(slide, "Our History", payload, primary_color)
    elif slide_type == "FootprintBlock":
        create_content_pptx_slide(slide, "Our Growing Footprint", payload, primary_color)
    elif slide_type == "RegionsBlock":
        create_content_pptx_slide(slide, "Our Regions", payload, primary_color)
    elif slide_type == "CultureBlock":
        create_culture_pptx_slide(slide, payload, brand, primary_color)
    elif slide_type == "GameSlide":
        create_game_pptx_slide(slide, payload, primary_color)
    elif slide_type == "PillarsClosing":
        create_pillars_pptx_slide(slide, payload, primary_color, secondary_color)
    else:
        # Generic slide
        create_generic_pptx_slide(slide, slide_type, payload, primary_color)


def create_welcome_pptx_slide(
    slide,
    payload: Dict[str, Any],
    brand: models.Brand,
    facility: models.Facility,
    primary_color: RGBColor
):
    """Create welcome/intro slide"""
    # Background color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"Welcome to {brand.name}"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Facility name
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4), Inches(12.33), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = facility.name
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Event type
    event_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(12.33), Inches(0.75)
    )
    event_frame = event_box.text_frame
    event_para = event_frame.paragraphs[0]
    event_para.text = "Culture Night"
    event_para.font.size = Pt(24)
    event_para.font.color.rgb = RGBColor(255, 255, 255)
    event_para.alignment = PP_ALIGN.CENTER


def create_raffle_pptx_slide(
    slide,
    payload: Dict[str, Any],
    primary_color: RGBColor
):
    """Create raffle bumper slide"""
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = payload.get("title", "RAFFLE TIME!")
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(12.33), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = payload.get("subtitle", "Time to draw a winner!")
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = primary_color
    subtitle_para.alignment = PP_ALIGN.CENTER


def create_content_pptx_slide(
    slide,
    title: str,
    payload: Dict[str, Any],
    primary_color: RGBColor
):
    """Create generic content slide"""
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = primary_color
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = payload.get("title", title)
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    content = payload.get("content", {})
    if isinstance(content, dict):
        # If content has items, display them
        items = content.get("items", [])
        if items:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(2), Inches(11.83), Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items[:6]):  # Limit to 6 items
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if isinstance(item, dict):
                    p.text = f"  {item.get('text', str(item))}"
                else:
                    p.text = f"  {item}"
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.space_after = Pt(12)


def create_culture_pptx_slide(
    slide,
    payload: Dict[str, Any],
    brand: models.Brand,
    primary_color: RGBColor
):
    """Create culture slide with Cascadia Way vs Common Way"""
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = primary_color
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = payload.get("title", f"The {brand.name} Way")
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.33), Inches(0.75)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "We are NOT corporate"
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = primary_color
    subtitle_para.alignment = PP_ALIGN.CENTER


def create_game_pptx_slide(
    slide,
    payload: Dict[str, Any],
    primary_color: RGBColor
):
    """Create game slide"""
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.5)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = primary_color
    title_bar.line.fill.background()

    # Game title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = payload.get("title", "Game Time!")
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Value label (if exists)
    value_label = payload.get("value_label")
    if value_label:
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12.33), Inches(0.6)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = value_label
        label_para.font.size = Pt(24)
        label_para.font.bold = True
        label_para.font.color.rgb = primary_color
        label_para.alignment = PP_ALIGN.CENTER

    # Rules
    rules = payload.get("rules", "")
    if rules:
        rules_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3), Inches(11.83), Inches(4)
        )
        rules_frame = rules_box.text_frame
        rules_frame.word_wrap = True
        rules_para = rules_frame.paragraphs[0]
        rules_para.text = rules
        rules_para.font.size = Pt(22)
        rules_para.font.color.rgb = RGBColor(50, 50, 50)


def create_pillars_pptx_slide(
    slide,
    payload: Dict[str, Any],
    primary_color: RGBColor,
    secondary_color: RGBColor
):
    """Create 3 pillars closing slide"""
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.75), Inches(12.33), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Our 3 Pillars"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Pillars
    pillars = payload.get("pillars", [
        {"name": "Clinical", "description": "Excellence in care"},
        {"name": "Cultural", "description": "Our people, our values"},
        {"name": "Financial", "description": "Sustainable success"},
    ])

    x_positions = [Inches(1.5), Inches(5.17), Inches(8.83)]
    for i, pillar in enumerate(pillars[:3]):
        # Pillar box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_positions[i], Inches(2.5), Inches(3), Inches(4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.fill.background()

        # Pillar name
        name_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), Inches(3), Inches(2.8), Inches(1)
        )
        name_frame = name_box.text_frame
        name_para = name_frame.paragraphs[0]
        name_para.text = pillar.get("name", "")
        name_para.font.size = Pt(28)
        name_para.font.bold = True
        name_para.font.color.rgb = primary_color
        name_para.alignment = PP_ALIGN.CENTER

        # Pillar description
        desc_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), Inches(4), Inches(2.8), Inches(2)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = pillar.get("description", "")
        desc_para.font.size = Pt(18)
        desc_para.font.color.rgb = RGBColor(100, 100, 100)
        desc_para.alignment = PP_ALIGN.CENTER


def create_generic_pptx_slide(
    slide,
    title: str,
    payload: Dict[str, Any],
    primary_color: RGBColor
):
    """Create a generic slide for unknown types"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.33), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color


def export_to_pdf(
    db: Session,
    presentation_id: int,
    output_path: Optional[str] = None
) -> bytes:
    """
    Export a presentation to PDF format.

    This creates a simple PDF from the slides. For production use,
    consider using a headless browser to render HTML slides to PDF.

    Args:
        db: Database session
        presentation_id: ID of the presentation to export
        output_path: Optional path to save the file

    Returns:
        Bytes of the PDF file
    """
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor, white, black

    # Get presentation data
    presentation = crud.get_presentation(db, presentation_id)
    if not presentation:
        raise ValueError(f"Presentation {presentation_id} not found")

    brand = crud.get_brand(db, presentation.brand_id)
    facility = crud.get_facility(db, presentation.facility_id)
    slides = crud.get_presentation_slides(db, presentation_id)

    # Create PDF
    buffer = io.BytesIO()
    page_size = landscape(A4)
    c = canvas.Canvas(buffer, pagesize=page_size)
    width, height = page_size

    primary_color = HexColor(brand.primary_color or "#0b7280")
    secondary_color = HexColor(brand.secondary_color or "#065a67")

    for slide_data in slides:
        add_slide_to_pdf(c, slide_data, brand, facility, width, height, primary_color)
        c.showPage()

    c.save()
    buffer.seek(0)
    pdf_bytes = buffer.read()

    if output_path:
        with open(output_path, 'wb') as f:
            f.write(pdf_bytes)

    return pdf_bytes


def add_slide_to_pdf(
    c,
    slide_data: models.PresentationSlide,
    brand: models.Brand,
    facility: models.Facility,
    width: float,
    height: float,
    primary_color
):
    """Add a single slide to the PDF"""
    from reportlab.lib.colors import white, black, HexColor

    slide_type = slide_data.slide_type
    payload = slide_data.payload or {}

    if slide_type == "WelcomeIntro":
        # Background
        c.setFillColor(primary_color)
        c.rect(0, 0, width, height, fill=1, stroke=0)
        # Title
        c.setFillColor(white)
        c.setFont("Helvetica-Bold", 48)
        c.drawCentredString(width/2, height - 200, f"Welcome to {brand.name}")
        # Facility
        c.setFont("Helvetica", 28)
        c.drawCentredString(width/2, height - 280, facility.name)
        # Event
        c.setFont("Helvetica", 22)
        c.drawCentredString(width/2, height - 340, "Culture Night")

    elif slide_type == "RaffleBumper":
        # Gold background
        c.setFillColor(HexColor("#FFD700"))
        c.rect(0, 0, width, height, fill=1, stroke=0)
        # Title
        c.setFillColor(primary_color)
        c.setFont("Helvetica-Bold", 64)
        c.drawCentredString(width/2, height - 250, payload.get("title", "RAFFLE TIME!"))
        c.setFont("Helvetica", 24)
        c.drawCentredString(width/2, height - 320, payload.get("subtitle", ""))

    elif slide_type == "GameSlide":
        # Header
        c.setFillColor(primary_color)
        c.rect(0, height - 80, width, 80, fill=1, stroke=0)
        c.setFillColor(white)
        c.setFont("Helvetica-Bold", 32)
        c.drawString(30, height - 55, payload.get("title", "Game Time!"))
        # Value label
        if payload.get("value_label"):
            c.setFillColor(primary_color)
            c.setFont("Helvetica-Bold", 24)
            c.drawCentredString(width/2, height - 130, payload.get("value_label"))
        # Rules
        c.setFillColor(black)
        c.setFont("Helvetica", 18)
        rules = payload.get("rules", "")
        y = height - 200
        for line in rules.split('\n')[:8]:
            c.drawString(50, y, line[:80])
            y -= 25

    elif slide_type == "PillarsClosing":
        # Background
        c.setFillColor(primary_color)
        c.rect(0, 0, width, height, fill=1, stroke=0)
        # Title
        c.setFillColor(white)
        c.setFont("Helvetica-Bold", 42)
        c.drawCentredString(width/2, height - 80, "Our 3 Pillars")
        # Pillars
        pillars = payload.get("pillars", [
            {"name": "Clinical"}, {"name": "Cultural"}, {"name": "Financial"}
        ])
        x_positions = [width/4 - 50, width/2, 3*width/4 + 50]
        for i, pillar in enumerate(pillars[:3]):
            # Box
            c.setFillColor(white)
            c.roundRect(x_positions[i] - 100, 150, 200, 280, 10, fill=1, stroke=0)
            # Name
            c.setFillColor(primary_color)
            c.setFont("Helvetica-Bold", 24)
            c.drawCentredString(x_positions[i], 360, pillar.get("name", ""))
            # Description
            c.setFillColor(black)
            c.setFont("Helvetica", 14)
            desc = pillar.get("description", "")
            c.drawCentredString(x_positions[i], 310, desc[:30])

    else:
        # Generic slide
        # Header
        c.setFillColor(primary_color)
        c.rect(0, height - 80, width, 80, fill=1, stroke=0)
        c.setFillColor(white)
        c.setFont("Helvetica-Bold", 32)
        title = payload.get("title", slide_type.replace("Block", ""))
        c.drawString(30, height - 55, title)


def get_export_filename(presentation: models.Presentation, extension: str) -> str:
    """Generate a filename for export"""
    safe_title = "".join(c for c in presentation.title if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_title = safe_title.replace(' ', '_')[:50]
    timestamp = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
    return f"{safe_title}_{timestamp}.{extension}"
